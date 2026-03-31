"""PowerPoint report generator using python-pptx."""
import io
import logging
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

COBALT = RGBColor(0x00, 0x2C, 0x77)
TEAL = RGBColor(0x00, 0xA8, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class PPTXGenerator:
    def __init__(self, template_path: Optional[Path] = None) -> None:
        self._template_path = template_path

    async def generate(self, context: Any) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        self._add_title_slide(prs, context)
        # Slide 2: KPI Summary
        self._add_kpi_slide(prs, context)
        # Slide 3: Top Variances
        self._add_variance_slide(prs, context)
        # Slide 4: Risk Items
        self._add_risk_slide(prs, context)
        # Slide 5: Summary
        self._add_summary_slide(prs, context)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _add_title_slide(self, prs, context):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COBALT

        # Title text
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Variance Analysis Report"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"{context.period_id} | {context.view} vs {context.comparison_base}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEAL
        p2.alignment = PP_ALIGN.CENTER

    def _add_kpi_slide(self, prs, context):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "Key Performance Indicators"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.color.rgb = COBALT

        # KPI boxes
        for i, card in enumerate(context.summary_cards[:5]):
            left = Inches(0.5 + i * 2.5)
            shape = slide.shapes.add_textbox(left, Inches(1.5), Inches(2.2), Inches(1.5))
            tf = shape.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = card.get("metric_name", "")
            p1.font.size = Pt(10)
            p1.font.color.rgb = TEAL
            p2 = tf.add_paragraph()
            p2.text = f"${card.get('actual', 0):,.0f}"
            p2.font.size = Pt(20)
            p2.font.bold = True
            p3 = tf.add_paragraph()
            pct = card.get("variance_pct", 0)
            p3.text = f"{pct:+.1f}%"
            p3.font.size = Pt(12)
            p3.font.color.rgb = RGBColor(0x2D, 0xD4, 0xA8) if pct >= 0 else RGBColor(0xF9, 0x70, 0x66)

    def _add_variance_slide(self, prs, context):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "Top Material Variances"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.color.rgb = COBALT

        sorted_vars = sorted(context.variances, key=lambda v: abs(v.get("variance_amount", 0)), reverse=True)[:8]
        rows = len(sorted_vars) + 1
        table = slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.3), Inches(12), Inches(rows * 0.5)).table
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(2)
        table.columns[3].width = Inches(4)

        for j, header in enumerate(["Account", "Variance $", "%", "Narrative"]):
            cell = table.cell(0, j)
            cell.text = header
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True

        for i, v in enumerate(sorted_vars):
            table.cell(i+1, 0).text = v.get("account_name", "")[:35]
            table.cell(i+1, 1).text = f"${v.get('variance_amount', 0):,.0f}"
            table.cell(i+1, 2).text = f"{v.get('variance_pct', 0):.1f}%"
            table.cell(i+1, 3).text = v.get("narrative_oneliner", "")[:50]

    def _add_risk_slide(self, prs, context):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "Risk Items & Alerts"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.color.rgb = COBALT

        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        for alert in (context.netting_alerts or [])[:3]:
            p = tf.add_paragraph()
            p.text = f"Netting: {alert.get('left', '')} / {alert.get('right', '')} (Net: {alert.get('net', '')})"
            p.font.size = Pt(12)
        for alert in (context.trend_alerts or [])[:3]:
            p = tf.add_paragraph()
            p.text = f"Trend: {alert.get('description', '')} - {alert.get('projection', '')}"
            p.font.size = Pt(12)

    def _add_summary_slide(self, prs, context):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        body = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = body.text_frame
        p = tf.paragraphs[0]
        p.text = "Marsh Vantage - Confidential"
        p.font.size = Pt(14)
        p.font.color.rgb = TEAL
        p.alignment = PP_ALIGN.CENTER
